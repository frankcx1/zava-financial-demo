"""Create the Stop Waiting Start Building PowerPoint slide - v2 larger layout."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Create presentation with widescreen dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Use blank layout
blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

# Colors
DARK_BG = RGBColor(15, 23, 42)
WHITE = RGBColor(255, 255, 255)
ACCENT_GREEN = RGBColor(34, 197, 94)
ACCENT_BLUE = RGBColor(59, 130, 246)
ACCENT_PURPLE = RGBColor(139, 92, 246)
LIGHT_GRAY = RGBColor(148, 163, 184)
MID_GRAY = RGBColor(100, 116, 139)
TEXT_LIGHT = RGBColor(226, 232, 240)

# Background
background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
background.fill.solid()
background.fill.fore_color.rgb = DARK_BG
background.line.fill.background()

# === TITLE ===
title = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(0.7))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "Stop Waiting. Start Building."
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = WHITE

subtitle = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(12), Inches(0.5))
tf = subtitle.text_frame
p = tf.paragraphs[0]
p.text = "The NPU Opportunity Is Now — And It's Yours to Create"
p.font.size = Pt(22)
p.font.color.rgb = LIGHT_GRAY

# === LEFT COLUMN: VIBE CODING + 80/20 ===

# VIBE CODING MOMENT
sec1_title = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(5.5), Inches(0.45))
tf = sec1_title.text_frame
p = tf.paragraphs[0]
p.text = "THE VIBE CODING MOMENT"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = ACCENT_GREEN

sec1_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.25), Inches(5.5), Inches(1.5))
tf = sec1_box.text_frame
tf.word_wrap = True

bullets = [
    "Claude Code + GitHub Copilot CLI = anyone can build production demos",
    "No SDK expertise required — describe what you want, iterate in real-time",
    "The best demos aren't coming from product teams — they're coming from you"
]
for i, bullet in enumerate(bullets):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = "• " + bullet
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_LIGHT
    p.space_after = Pt(10)

# 80/20 RULE BOX
rule_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(4.0), Inches(5.5), Inches(1.6))
rule_box.fill.solid()
rule_box.fill.fore_color.rgb = RGBColor(30, 41, 59)
rule_box.line.color.rgb = RGBColor(71, 85, 105)
rule_box.line.width = Pt(1)

rule_title = slide.shapes.add_textbox(Inches(0.85), Inches(4.15), Inches(4), Inches(0.4))
tf = rule_title.text_frame
p = tf.paragraphs[0]
p.text = "THE 80/20 RULE"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = WHITE

rule_content = slide.shapes.add_textbox(Inches(0.85), Inches(4.55), Inches(5.1), Inches(1.0))
tf = rule_content.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "80% Local"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = ACCENT_GREEN
run = p.add_run()
run.text = " — routine tasks, sensitive data, always-on automation "
run.font.size = Pt(14)
run.font.color.rgb = RGBColor(203, 213, 225)
run = p.add_run()
run.text = "(free)"
run.font.size = Pt(14)
run.font.color.rgb = ACCENT_GREEN

p = tf.add_paragraph()
p.space_before = Pt(8)
p.text = "20% Frontier"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE
run = p.add_run()
run.text = " — complex reasoning, novel problems, maximum capability"
run.font.size = Pt(14)
run.font.color.rgb = RGBColor(203, 213, 225)

# === MIDDLE COLUMN: TWO-BRAIN STRATEGY ===

sec2_title = slide.shapes.add_textbox(Inches(6.4), Inches(1.8), Inches(3.5), Inches(0.45))
tf = sec2_title.text_frame
p = tf.paragraphs[0]
p.text = "THE TWO-BRAIN STRATEGY"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE

# Table
headers = ["", "Cloud AI", "Local AI (NPU)"]
rows = [
    ["Role", "Expert Consultant", "Chief of Staff"],
    ["When", "Big projects", "Always at your desk"],
    ["Strength", "Brilliant problem-solver", "Knows your context"],
    ["Data", "Leaves the building", "Never leaves"],
]

table_top = Inches(2.25)
row_height = Inches(0.38)
col_widths = [Inches(0.95), Inches(1.55), Inches(1.55)]

for row_idx, row_data in enumerate([headers] + rows):
    y = table_top + row_idx * row_height
    for col_idx, cell_text in enumerate(row_data):
        x = Inches(6.4) + sum(col_widths[:col_idx], Inches(0))
        cell = slide.shapes.add_textbox(x, y, col_widths[col_idx], row_height)
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = cell_text
        p.font.size = Pt(12)
        if row_idx == 0:
            p.font.bold = True
            p.font.color.rgb = LIGHT_GRAY
        elif col_idx == 0:
            p.font.bold = True
            p.font.color.rgb = MID_GRAY
        elif col_idx == 1:
            p.font.color.rgb = RGBColor(147, 197, 253)
        else:
            p.font.color.rgb = RGBColor(134, 239, 172)

# Quote
quote = slide.shapes.add_textbox(Inches(6.4), Inches(4.25), Inches(4), Inches(0.6))
tf = quote.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '"Biggest brain for development.\nMost private brain for production."'
p.font.size = Pt(13)
p.font.italic = True
p.font.color.rgb = RGBColor(203, 213, 225)

# === RIGHT COLUMN: TOKENOMICS ===

sec3_title = slide.shapes.add_textbox(Inches(10.5), Inches(1.8), Inches(2.5), Inches(0.45))
tf = sec3_title.text_frame
p = tf.paragraphs[0]
p.text = "TOKENOMICS"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = ACCENT_PURPLE

sec3_box = slide.shapes.add_textbox(Inches(10.5), Inches(2.25), Inches(2.5), Inches(2.5))
tf = sec3_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Subscription AI"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = RGBColor(196, 181, 253)

p = tf.add_paragraph()
p.text = "Fixed cost, human-speed"
p.font.size = Pt(11)
p.font.color.rgb = LIGHT_GRAY
p.space_after = Pt(12)

p = tf.add_paragraph()
p.text = "Token-based AI (API)"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = RGBColor(196, 181, 253)

p = tf.add_paragraph()
p.text = "Computer-to-computer burns tokens fast:"
p.font.size = Pt(11)
p.font.color.rgb = LIGHT_GRAY
p.space_after = Pt(4)

for item in ["• 100s of agent calls/day", "• 24/7 automation", "• Costs compound quietly"]:
    p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(203, 213, 225)
    p.space_after = Pt(3)

# === CALL TO ACTION (full width bottom) ===
cta_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.85), Inches(12.1), Inches(1.25))
cta_box.fill.solid()
cta_box.fill.fore_color.rgb = RGBColor(20, 45, 35)
cta_box.line.color.rgb = ACCENT_GREEN
cta_box.line.width = Pt(2)

cta_title = slide.shapes.add_textbox(Inches(0.9), Inches(5.95), Inches(3), Inches(0.4))
tf = cta_title.text_frame
p = tf.paragraphs[0]
p.text = "CALL TO ACTION"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = ACCENT_GREEN

cta_content = slide.shapes.add_textbox(Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.7))
tf = cta_content.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "NPU is in every new Surface and across Windows — Copilot+ PCs are the new baseline.  "
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(187, 247, 208)

run = p.add_run()
run.text = "The use cases won't discover themselves.  "
run.font.size = Pt(14)
run.font.color.rgb = RGBColor(187, 247, 208)

run = p.add_run()
run.text = "Build the demo. Show the value. Lead the conversation."
run.font.size = Pt(14)
run.font.bold = True
run.font.color.rgb = WHITE

# === FOOTER ===
footer = slide.shapes.add_textbox(Inches(10.5), Inches(7.15), Inches(2.5), Inches(0.25))
tf = footer.text_frame
p = tf.paragraphs[0]
p.text = "Microsoft Surface | 2026"
p.font.size = Pt(9)
p.font.color.rgb = MID_GRAY
p.alignment = PP_ALIGN.RIGHT

# Save
output_path = r"C:\Users\frankbu\OneDrive - Microsoft\NPU\Stop_Waiting_Start_Building_v2.pptx"
prs.save(output_path)
print(f"Saved to: {output_path}")
